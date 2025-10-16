#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Detector e Reparador Universal de Arquivos
Script completo para detectar e reparar problemas em arquivos
Combina todas as funcionalidades em uma única ferramenta
"""

import os
import sys
import json
import shutil
import hashlib
import logging
import zipfile
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Set
import argparse

# Bibliotecas para manipulação de arquivos
try:
    import pandas as pd
    from docx import Document
    import PyPDF2
    from pptx import Presentation
    from openpyxl import load_workbook, Workbook
    from openpyxl.utils.exceptions import InvalidFileException
    import fitz  # PyMuPDF
    import docx2txt
except ImportError as e:
    print(f"Erro ao importar bibliotecas: {e}")
    print("Execute: pip install pandas python-docx PyPDF2 python-pptx openpyxl PyMuPDF docx2txt")
    sys.exit(1)

class DetectorReparadorUniversal:
    def __init__(self, diretorio_base: str, modo: str = "completo"):
        """
        Inicializa o detector/reparador universal
        
        Args:
            diretorio_base: Diretório base para análise
            modo: 'deteccao', 'reparo' ou 'completo'
        """
        self.diretorio_base = Path(diretorio_base)
        self.diretorio_consolidado = self.diretorio_base / "Consolidado"
        self.diretorio_logs = self.diretorio_consolidado / "logs"
        self.modo = modo
        
        # Criar diretórios necessários
        self.diretorio_logs.mkdir(exist_ok=True)
        
        if modo in ['reparo', 'completo']:
            self.diretorio_reparados = self.diretorio_consolidado / "arquivos_reparados"
            self.diretorio_backup = self.diretorio_consolidado / "backup_originais"
            self.diretorio_reparados.mkdir(exist_ok=True)
            self.diretorio_backup.mkdir(exist_ok=True)
        
        # Configurar logging
        self.timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        log_file = self.diretorio_logs / f"detector_reparador_universal_{self.timestamp}.log"
        
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(levelname)s - %(message)s',
            handlers=[
                logging.FileHandler(log_file, encoding='utf-8'),
                logging.StreamHandler()
            ]
        )
        self.logger = logging.getLogger(__name__)
        
        # Estatísticas
        self.stats = {
            'arquivos_analisados': 0,
            'arquivos_processados': 0,
            'arquivos_reparados': 0,
            'duplicados_removidos': 0,
            'extensoes_corrigidas': 0,
            'pdfs_reparados': 0,
            'xlsx_reparados': 0,
            'docx_reparados': 0,
            'pptx_reparados': 0,
            'problemas_encontrados': 0,
            'erros': 0
        }
        
        # Resultados
        self.problemas_detectados = {
            'arquivos_vazios': [],
            'arquivos_muito_pequenos': [],
            'arquivos_com_conteudo_suspeito': [],
            'arquivos_com_estrutura_danificada': [],
            'arquivos_com_metadados_corrompidos': [],
            'arquivos_com_encoding_problemas': [],
            'arquivos_duplicados': [],
            'arquivos_com_extensao_incorreta': []
        }
        
        self.problemas_corrigidos = []
        self.erros_encontrados = []
        self.checksums_calculados = {}

    def calcular_checksum(self, arquivo: Path) -> str:
        """Calcula o checksum MD5 de um arquivo"""
        if str(arquivo) in self.checksums_calculados:
            return self.checksums_calculados[str(arquivo)]
        
        hash_md5 = hashlib.md5()
        try:
            with open(arquivo, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_md5.update(chunk)
            checksum = hash_md5.hexdigest()
            self.checksums_calculados[str(arquivo)] = checksum
            return checksum
        except Exception as e:
            self.logger.error(f"Erro ao calcular checksum de {arquivo}: {e}")
            return ""

    def detectar_tipo_arquivo(self, arquivo: Path) -> Optional[str]:
        """Detecta o tipo real do arquivo baseado no conteúdo"""
        try:
            with open(arquivo, 'rb') as f:
                header = f.read(8)
            
            # Assinaturas de arquivos
            assinaturas = {
                b'\x50\x4B\x03\x04': '.zip',  # ZIP/Office
                b'\x50\x4B\x05\x06': '.zip',  # ZIP vazio
                b'\x50\x4B\x07\x08': '.zip',  # ZIP spanned
                b'\x25\x50\x44\x46': '.pdf',  # PDF
                b'\xFF\xD8\xFF': '.jpg',       # JPEG
                b'\x89\x50\x4E\x47': '.png',   # PNG
                b'\x47\x49\x46\x38': '.gif',   # GIF
            }
            
            for assinatura, extensao in assinaturas.items():
                if header.startswith(assinatura):
                    if extensao == '.zip':
                        return self._detectar_tipo_office(arquivo)
                    return extensao
            
            return None
            
        except Exception:
            return None

    def _detectar_tipo_office(self, arquivo: Path) -> Optional[str]:
        """Detecta tipo específico de arquivo Office"""
        try:
            with zipfile.ZipFile(arquivo, 'r') as zip_ref:
                files = zip_ref.namelist()
                
                if 'word/document.xml' in files:
                    return '.docx'
                elif 'xl/workbook.xml' in files:
                    return '.xlsx'
                elif 'ppt/presentation.xml' in files:
                    return '.pptx'
                else:
                    return '.zip'
                    
        except Exception:
            return '.zip'

    def verificar_arquivo_vazio(self, arquivo: Path) -> bool:
        """Verifica se o arquivo está vazio"""
        try:
            return arquivo.stat().st_size == 0
        except Exception:
            return False

    def verificar_arquivo_muito_pequeno(self, arquivo: Path, limite: int = 100) -> bool:
        """Verifica se o arquivo é suspeitosamente pequeno"""
        try:
            tamanho = arquivo.stat().st_size
            extensao = arquivo.suffix.lower()
            
            # Limites específicos por tipo
            limites = {
                '.docx': 1000,
                '.xlsx': 2000,
                '.pptx': 5000,
                '.pdf': 500,
                '.txt': 10
            }
            
            limite_especifico = limites.get(extensao, limite)
            return tamanho < limite_especifico
        except Exception:
            return False

    def verificar_conteudo_suspeito(self, arquivo: Path) -> Optional[Dict]:
        """Verifica se o arquivo tem conteúdo suspeito"""
        extensao = arquivo.suffix.lower()
        
        try:
            if extensao == '.pdf':
                return self._verificar_pdf_suspeito(arquivo)
            elif extensao == '.docx':
                return self._verificar_docx_suspeito(arquivo)
            elif extensao == '.xlsx':
                return self._verificar_xlsx_suspeito(arquivo)
            elif extensao == '.pptx':
                return self._verificar_pptx_suspeito(arquivo)
        except Exception as e:
            return {
                'problema': f'Erro ao verificar {extensao.upper()}',
                'detalhes': str(e)
            }
        
        return None

    def _verificar_pdf_suspeito(self, arquivo: Path) -> Optional[Dict]:
        """Verifica problemas específicos em PDFs"""
        try:
            # Tentar com PyMuPDF
            doc = fitz.open(arquivo)
            texto_total = ""
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                texto_total += page.get_text()
            
            doc.close()
            
            if not texto_total.strip():
                return {
                    'problema': 'PDF sem texto extraível',
                    'detalhes': 'Todas as páginas estão sem texto ou são imagens'
                }
                
        except Exception as e:
            return {
                'problema': 'Erro ao ler PDF',
                'detalhes': str(e)
            }
        
        return None

    def _verificar_docx_suspeito(self, arquivo: Path) -> Optional[Dict]:
        """Verifica problemas específicos em DOCX"""
        try:
            doc = Document(arquivo)
            texto_total = ""
            
            for paragrafo in doc.paragraphs:
                texto_total += paragrafo.text
            
            if not texto_total.strip():
                return {
                    'problema': 'DOCX sem conteúdo de texto',
                    'detalhes': 'Documento não contém texto legível'
                }
                
        except Exception as e:
            return {
                'problema': 'Erro ao ler DOCX',
                'detalhes': str(e)
            }
        
        return None

    def _verificar_xlsx_suspeito(self, arquivo: Path) -> Optional[Dict]:
        """Verifica problemas específicos em XLSX"""
        try:
            wb = load_workbook(arquivo, data_only=True)
            tem_dados = False
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value is not None:
                            tem_dados = True
                            break
                    if tem_dados:
                        break
                if tem_dados:
                    break
            
            wb.close()
            
            if not tem_dados:
                return {
                    'problema': 'XLSX sem dados',
                    'detalhes': 'Planilha não contém dados nas células'
                }
                
        except Exception as e:
            return {
                'problema': 'Erro ao ler XLSX',
                'detalhes': str(e)
            }
        
        return None

    def _verificar_pptx_suspeito(self, arquivo: Path) -> Optional[Dict]:
        """Verifica problemas específicos em PPTX"""
        try:
            prs = Presentation(arquivo)
            
            if len(prs.slides) == 0:
                return {
                    'problema': 'PPTX sem slides',
                    'detalhes': 'Apresentação não contém slides'
                }
                
        except Exception as e:
            return {
                'problema': 'Erro ao ler PPTX',
                'detalhes': str(e)
            }
        
        return None

    def detectar_duplicados(self, arquivos: List[Path]) -> List[Dict]:
        """Detecta arquivos duplicados baseado em checksum"""
        checksums = {}
        duplicados = []
        
        for arquivo in arquivos:
            if not arquivo.is_file():
                continue
                
            checksum = self.calcular_checksum(arquivo)
            if not checksum:
                continue
            
            if checksum in checksums:
                duplicados.append({
                    'arquivo': str(arquivo),
                    'problema': 'Arquivo duplicado',
                    'detalhes': f'Idêntico a: {checksums[checksum]}',
                    'checksum': checksum
                })
            else:
                checksums[checksum] = str(arquivo)
        
        return duplicados

    def detectar_extensoes_incorretas(self, arquivos: List[Path]) -> List[Dict]:
        """Detecta arquivos com extensões incorretas"""
        incorretas = []
        
        for arquivo in arquivos:
            if not arquivo.is_file():
                continue
            
            tipo_real = self.detectar_tipo_arquivo(arquivo)
            extensao_atual = arquivo.suffix.lower()
            
            if tipo_real and tipo_real != extensao_atual:
                incorretas.append({
                    'arquivo': str(arquivo),
                    'problema': 'Extensão incorreta',
                    'detalhes': f'Tipo real: {tipo_real}, Extensão atual: {extensao_atual}'
                })
        
        return incorretas

    def executar_deteccao(self) -> Dict:
        """Executa a detecção completa de problemas"""
        self.logger.info("Iniciando detecção de problemas...")
        
        # Obter lista de arquivos
        arquivos = []
        for arquivo in self.diretorio_consolidado.rglob('*'):
            if arquivo.is_file() and not arquivo.name.startswith('.'):
                # Ignorar arquivos em subpastas de processamento
                if not any(parte in str(arquivo) for parte in ['logs', 'arquivos_txt', 'arquivos_csv', 'backup', 'reparados']):
                    arquivos.append(arquivo)
        
        self.stats['arquivos_analisados'] = len(arquivos)
        self.logger.info(f"Analisando {len(arquivos)} arquivos...")
        
        # Detectar problemas
        for arquivo in arquivos:
            try:
                # Arquivo vazio
                if self.verificar_arquivo_vazio(arquivo):
                    self.problemas_detectados['arquivos_vazios'].append({
                        'arquivo': str(arquivo),
                        'problema': 'Arquivo vazio',
                        'detalhes': 'Arquivo com 0 bytes'
                    })
                
                # Arquivo muito pequeno
                elif self.verificar_arquivo_muito_pequeno(arquivo):
                    self.problemas_detectados['arquivos_muito_pequenos'].append({
                        'arquivo': str(arquivo),
                        'problema': 'Arquivo suspeitosamente pequeno',
                        'detalhes': f'Tamanho: {arquivo.stat().st_size} bytes'
                    })
                
                # Conteúdo suspeito
                problema_conteudo = self.verificar_conteudo_suspeito(arquivo)
                if problema_conteudo:
                    self.problemas_detectados['arquivos_com_conteudo_suspeito'].append({
                        'arquivo': str(arquivo),
                        **problema_conteudo
                    })
                
            except Exception as e:
                self.logger.error(f"Erro ao analisar {arquivo}: {e}")
                self.erros_encontrados.append({
                    'arquivo': str(arquivo),
                    'erro': str(e),
                    'tipo': 'Análise'
                })
        
        # Detectar duplicados
        self.problemas_detectados['arquivos_duplicados'] = self.detectar_duplicados(arquivos)
        
        # Detectar extensões incorretas
        self.problemas_detectados['arquivos_com_extensao_incorreta'] = self.detectar_extensoes_incorretas(arquivos)
        
        # Calcular total de problemas
        total_problemas = sum(len(lista) for lista in self.problemas_detectados.values())
        self.stats['problemas_encontrados'] = total_problemas
        
        self.logger.info(f"Detecção concluída. {total_problemas} problemas encontrados.")
        
        return self.problemas_detectados

    def criar_backup(self, arquivo: Path) -> bool:
        """Cria backup do arquivo original"""
        try:
            backup_path = self.diretorio_backup / arquivo.name
            # Evitar conflitos de nome
            contador = 1
            while backup_path.exists():
                nome_base = arquivo.stem
                extensao = arquivo.suffix
                backup_path = self.diretorio_backup / f"{nome_base}_backup_{contador}{extensao}"
                contador += 1
            
            shutil.copy2(arquivo, backup_path)
            self.logger.info(f"Backup criado: {backup_path}")
            return True
        except Exception as e:
            self.logger.error(f"Erro ao criar backup de {arquivo}: {e}")
            return False

    def reparar_arquivo(self, arquivo: Path, problema: Dict) -> bool:
        """Repara um arquivo específico baseado no problema detectado"""
        tipo_problema = problema.get('problema', '')
        
        try:
            if 'PDF sem texto extraível' in tipo_problema:
                return self._reparar_pdf(arquivo)
            elif 'Erro ao ler XLSX' in tipo_problema:
                return self._reparar_xlsx(arquivo)
            elif 'Erro ao ler DOCX' in tipo_problema:
                return self._reparar_docx(arquivo)
            elif 'Erro ao ler PPTX' in tipo_problema:
                return self._reparar_pptx(arquivo)
            elif 'Extensão incorreta' in tipo_problema:
                return self._corrigir_extensao(arquivo, problema)
            elif 'Arquivo duplicado' in tipo_problema:
                return self._remover_duplicado(arquivo, problema)
            else:
                self.logger.warning(f"Tipo de problema não suportado para reparo: {tipo_problema}")
                return False
                
        except Exception as e:
            self.logger.error(f"Erro ao reparar {arquivo}: {e}")
            self.erros_encontrados.append({
                'arquivo': str(arquivo),
                'erro': str(e),
                'tipo': 'Reparo'
            })
            return False

    def _reparar_pdf(self, arquivo: Path) -> bool:
        """Repara PDF com problemas"""
        try:
            if not self.criar_backup(arquivo):
                return False
            
            doc = fitz.open(arquivo)
            novo_arquivo = self.diretorio_reparados / arquivo.name
            novo_doc = fitz.open()
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                
                if not text.strip():
                    # Converter página para imagem e reinserir
                    pix = page.get_pixmap()
                    img_data = pix.tobytes("png")
                    new_page = novo_doc.new_page(width=page.rect.width, height=page.rect.height)
                    new_page.insert_image(page.rect, stream=img_data)
                else:
                    novo_doc.insert_pdf(doc, from_page=page_num, to_page=page_num)
            
            novo_doc.save(novo_arquivo)
            novo_doc.close()
            doc.close()
            
            self.stats['pdfs_reparados'] += 1
            self.problemas_corrigidos.append({
                'arquivo': str(arquivo),
                'problema': 'PDF sem texto extraível',
                'solucao': 'PDF recriado com estrutura corrigida',
                'arquivo_reparado': str(novo_arquivo)
            })
            return True
            
        except Exception as e:
            self.logger.error(f"Erro ao reparar PDF {arquivo}: {e}")
            return False

    def _reparar_xlsx(self, arquivo: Path) -> bool:
        """Repara XLSX com problemas"""
        try:
            if not self.criar_backup(arquivo):
                return False
            
            # Tentar diferentes métodos
            metodos = [
                lambda: self._reparar_xlsx_openpyxl(arquivo),
                lambda: self._reparar_xlsx_pandas(arquivo),
                lambda: self._reparar_xlsx_zip(arquivo)
            ]
            
            for metodo in metodos:
                if metodo():
                    self.stats['xlsx_reparados'] += 1
                    return True
            
            return False
            
        except Exception as e:
            self.logger.error(f"Erro ao reparar XLSX {arquivo}: {e}")
            return False

    def _reparar_xlsx_openpyxl(self, arquivo: Path) -> bool:
        """Repara XLSX usando openpyxl"""
        try:
            wb = load_workbook(arquivo, data_only=True)
            novo_arquivo = self.diretorio_reparados / arquivo.name
            wb.save(novo_arquivo)
            wb.close()
            
            self.problemas_corrigidos.append({
                'arquivo': str(arquivo),
                'problema': 'Erro ao ler XLSX',
                'solucao': 'Arquivo recriado com openpyxl',
                'arquivo_reparado': str(novo_arquivo)
            })
            return True
            
        except Exception:
            return False

    def _reparar_xlsx_pandas(self, arquivo: Path) -> bool:
        """Repara XLSX usando pandas"""
        try:
            df = pd.read_excel(arquivo, sheet_name=None)
            novo_arquivo = self.diretorio_reparados / arquivo.name
            
            with pd.ExcelWriter(novo_arquivo, engine='openpyxl') as writer:
                for sheet_name, sheet_df in df.items():
                    sheet_df.to_excel(writer, sheet_name=sheet_name, index=False)
            
            self.problemas_corrigidos.append({
                'arquivo': str(arquivo),
                'problema': 'Erro ao ler XLSX',
                'solucao': 'Arquivo recriado com pandas',
                'arquivo_reparado': str(novo_arquivo)
            })
            return True
            
        except Exception:
            return False

    def _reparar_xlsx_zip(self, arquivo: Path) -> bool:
        """Repara XLSX como arquivo ZIP"""
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir)
                
                with zipfile.ZipFile(arquivo, 'r') as zip_ref:
                    zip_ref.extractall(temp_path)
                
                novo_arquivo = self.diretorio_reparados / arquivo.name
                with zipfile.ZipFile(novo_arquivo, 'w', zipfile.ZIP_DEFLATED) as zip_ref:
                    for file_path in temp_path.rglob('*'):
                        if file_path.is_file():
                            arcname = file_path.relative_to(temp_path)
                            zip_ref.write(file_path, arcname)
                
                # Verificar se funciona
                wb = load_workbook(novo_arquivo)
                wb.close()
                
                self.problemas_corrigidos.append({
                    'arquivo': str(arquivo),
                    'problema': 'Erro ao ler XLSX',
                    'solucao': 'Arquivo recriado como ZIP',
                    'arquivo_reparado': str(novo_arquivo)
                })
                return True
                
        except Exception:
            return False

    def _reparar_docx(self, arquivo: Path) -> bool:
        """Repara DOCX com problemas"""
        try:
            if not self.criar_backup(arquivo):
                return False
            
            # Tentar extrair texto e recriar documento
            texto = docx2txt.process(arquivo)
            
            novo_doc = Document()
            novo_doc.add_paragraph(texto)
            
            novo_arquivo = self.diretorio_reparados / arquivo.name
            novo_doc.save(novo_arquivo)
            
            self.stats['docx_reparados'] += 1
            self.problemas_corrigidos.append({
                'arquivo': str(arquivo),
                'problema': 'Erro ao ler DOCX',
                'solucao': 'Documento recriado com texto extraído',
                'arquivo_reparado': str(novo_arquivo)
            })
            return True
            
        except Exception as e:
            self.logger.error(f"Erro ao reparar DOCX {arquivo}: {e}")
            return False

    def _reparar_pptx(self, arquivo: Path) -> bool:
        """Repara PPTX com problemas"""
        try:
            if not self.criar_backup(arquivo):
                return False
            
            # Tentar recriar apresentação básica
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            title.text = f"Arquivo reparado: {arquivo.name}"
            
            novo_arquivo = self.diretorio_reparados / arquivo.name
            prs.save(novo_arquivo)
            
            self.stats['pptx_reparados'] += 1
            self.problemas_corrigidos.append({
                'arquivo': str(arquivo),
                'problema': 'Erro ao ler PPTX',
                'solucao': 'Apresentação recriada',
                'arquivo_reparado': str(novo_arquivo)
            })
            return True
            
        except Exception as e:
            self.logger.error(f"Erro ao reparar PPTX {arquivo}: {e}")
            return False

    def _corrigir_extensao(self, arquivo: Path, problema: Dict) -> bool:
        """Corrige extensão incorreta"""
        try:
            if not self.criar_backup(arquivo):
                return False
            
            tipo_real = self.detectar_tipo_arquivo(arquivo)
            if not tipo_real:
                return False
            
            novo_nome = arquivo.stem + tipo_real
            novo_caminho = arquivo.parent / novo_nome
            
            # Evitar conflitos
            contador = 1
            while novo_caminho.exists():
                novo_nome = f"{arquivo.stem}_{contador}{tipo_real}"
                novo_caminho = arquivo.parent / novo_nome
                contador += 1
            
            arquivo.rename(novo_caminho)
            
            self.stats['extensoes_corrigidas'] += 1
            self.problemas_corrigidos.append({
                'arquivo': str(arquivo),
                'problema': 'Extensão incorreta',
                'solucao': f'Extensão corrigida para {tipo_real}',
                'arquivo_corrigido': str(novo_caminho)
            })
            return True
            
        except Exception as e:
            self.logger.error(f"Erro ao corrigir extensão de {arquivo}: {e}")
            return False

    def _remover_duplicado(self, arquivo: Path, problema: Dict) -> bool:
        """Remove arquivo duplicado"""
        try:
            if not self.criar_backup(arquivo):
                return False
            
            pasta_duplicados = self.diretorio_consolidado / "duplicados_removidos"
            pasta_duplicados.mkdir(exist_ok=True)
            
            destino = pasta_duplicados / arquivo.name
            contador = 1
            while destino.exists():
                nome_base = arquivo.stem
                extensao = arquivo.suffix
                destino = pasta_duplicados / f"{nome_base}_dup_{contador}{extensao}"
                contador += 1
            
            shutil.move(arquivo, destino)
            
            self.stats['duplicados_removidos'] += 1
            self.problemas_corrigidos.append({
                'arquivo': str(arquivo),
                'problema': 'Arquivo duplicado',
                'solucao': 'Arquivo movido para pasta de duplicados',
                'arquivo_movido': str(destino)
            })
            return True
            
        except Exception as e:
            self.logger.error(f"Erro ao remover duplicado {arquivo}: {e}")
            return False

    def executar_reparo(self, problemas: Optional[Dict] = None) -> Dict:
        """Executa o reparo dos problemas detectados"""
        if problemas is None:
            problemas = self.problemas_detectados
        
        self.logger.info("Iniciando processo de reparo...")
        
        total_reparos = 0
        
        # Reparar cada categoria de problema
        for categoria, lista_problemas in problemas.items():
            if not lista_problemas:
                continue
                
            self.logger.info(f"Reparando {len(lista_problemas)} problemas de: {categoria}")
            
            for problema in lista_problemas:
                arquivo = Path(problema['arquivo'])
                if arquivo.exists():
                    if self.reparar_arquivo(arquivo, problema):
                        total_reparos += 1
                        self.stats['arquivos_processados'] += 1
        
        self.stats['arquivos_reparados'] = total_reparos
        self.logger.info(f"Reparo concluído. {total_reparos} arquivos reparados.")
        
        return {
            'arquivos_reparados': total_reparos,
            'problemas_corrigidos': self.problemas_corrigidos,
            'erros': self.erros_encontrados
        }

    def gerar_relatorio(self, incluir_deteccao: bool = True, incluir_reparo: bool = True) -> str:
        """Gera relatório completo"""
        relatorio = {
            'timestamp': self.timestamp,
            'modo_execucao': self.modo,
            'estatisticas': self.stats
        }
        
        if incluir_deteccao:
            relatorio['problemas_detectados'] = self.problemas_detectados
        
        if incluir_reparo:
            relatorio['problemas_corrigidos'] = self.problemas_corrigidos
            relatorio['erros_encontrados'] = self.erros_encontrados
        
        if hasattr(self, 'diretorio_reparados'):
            relatorio['diretorios'] = {
                'arquivos_reparados': str(self.diretorio_reparados),
                'backup_originais': str(self.diretorio_backup),
                'logs': str(self.diretorio_logs)
            }
        
        arquivo_relatorio = self.diretorio_logs / f"relatorio_universal_{self.timestamp}.json"
        
        with open(arquivo_relatorio, 'w', encoding='utf-8') as f:
            json.dump(relatorio, f, indent=2, ensure_ascii=False)
        
        return str(arquivo_relatorio)

    def executar(self) -> str:
        """Executa o processo completo baseado no modo"""
        if self.modo in ['deteccao', 'completo']:
            self.executar_deteccao()
        
        if self.modo in ['reparo', 'completo']:
            self.executar_reparo()
        
        return self.gerar_relatorio(
            incluir_deteccao=(self.modo in ['deteccao', 'completo']),
            incluir_reparo=(self.modo in ['reparo', 'completo'])
        )

def main():
    parser = argparse.ArgumentParser(description='Detector e Reparador Universal de Arquivos')
    parser.add_argument('--diretorio', '-d', 
                       default="/Users/marcosdaniels/Downloads/5. FULL SALES 2",
                       help='Diretório base para análise')
    parser.add_argument('--modo', '-m', 
                       choices=['deteccao', 'reparo', 'completo'],
                       default='completo',
                       help='Modo de operação')
    
    args = parser.parse_args()
    
    print("🔧 DETECTOR E REPARADOR UNIVERSAL DE ARQUIVOS")
    print("=" * 60)
    print(f"📁 Diretório: {args.diretorio}")
    print(f"🔄 Modo: {args.modo}")
    print("=" * 60)
    
    # Inicializar detector/reparador
    detector_reparador = DetectorReparadorUniversal(args.diretorio, args.modo)
    
    # Executar processo
    relatorio_final = detector_reparador.executar()
    
    # Exibir resultados
    print("\n" + "=" * 60)
    print("RELATÓRIO FINAL")
    print("=" * 60)
    print(f"Arquivos analisados: {detector_reparador.stats['arquivos_analisados']}")
    print(f"Problemas encontrados: {detector_reparador.stats['problemas_encontrados']}")
    
    if args.modo in ['reparo', 'completo']:
        print(f"Arquivos processados: {detector_reparador.stats['arquivos_processados']}")
        print(f"Arquivos reparados: {detector_reparador.stats['arquivos_reparados']}")
        print(f"Duplicados removidos: {detector_reparador.stats['duplicados_removidos']}")
        print(f"Extensões corrigidas: {detector_reparador.stats['extensoes_corrigidas']}")
        print(f"PDFs reparados: {detector_reparador.stats['pdfs_reparados']}")
        print(f"XLSX reparados: {detector_reparador.stats['xlsx_reparados']}")
        print(f"DOCX reparados: {detector_reparador.stats['docx_reparados']}")
        print(f"PPTX reparados: {detector_reparador.stats['pptx_reparados']}")
    
    print(f"Erros encontrados: {detector_reparador.stats['erros']}")
    print("=" * 60)
    print(f"\n📊 Relatório completo: {relatorio_final}")
    
    if hasattr(detector_reparador, 'diretorio_reparados'):
        print(f"🔧 Arquivos reparados: {detector_reparador.diretorio_reparados}")
        print(f"💾 Backups: {detector_reparador.diretorio_backup}")
    
    print("\nProcesso concluído! ✅")

if __name__ == "__main__":
    main()