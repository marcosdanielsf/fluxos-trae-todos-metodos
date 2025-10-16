#!/usr/bin/env python3
"""
Script de Consolidação e Conversão de Arquivos
Reorganiza arquivos de múltiplas pastas em uma estrutura consolidada
e converte para formato TXT com verificação de qualidade.
"""

import os
import shutil
import hashlib
import logging
import json
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Tuple, Optional
import traceback

# Bibliotecas para conversão de arquivos
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class ConsolidadorArquivos:
    """Classe principal para consolidação e conversão de arquivos."""
    
    def __init__(self, diretorio_base: str):
        self.diretorio_base = Path(diretorio_base)
        self.diretorio_consolidado = self.diretorio_base / "Consolidado"
        self.diretorio_logs = self.diretorio_consolidado / "logs"
        self.diretorio_txt = self.diretorio_consolidado / "arquivos_txt"
        
        # Mapeamento de categorias
        self.categorias = {
            "1.  Apresentações comerciais": "P1. Apresentações Comerciais",
            "2.  Apresentações Pitch em Eventos Presenciais": "P2. Pitch Eventos",
            "3. Produtos e Ofertas": "P3. Produtos e Ofertas",
            "4. Processos de RH Comercial  ": "P4. RH Comercial",
            "5. Pré Vendas  - FSS": "P5. Pré Vendas",
            "6. Vendas - FSS": "P6. Vendas",
            "7. Social Selling": "P7. Social Selling",
            "8. Sales Farm": "P8. Sales Farm",
            "9. Prospecção Fria": "P9. Prospecção Fria",
            "10. Programa de Embaixadores": "P10. Embaixadores",
            "11. Arquivos de Gestão": "P11. Gestão"
        }
        
        # Estatísticas do processo
        self.estatisticas = {
            "total_arquivos": 0,
            "arquivos_processados": 0,
            "arquivos_com_erro": 0,
            "conversoes_sucesso": 0,
            "conversoes_erro": 0,
            "formatos_encontrados": {},
            "erros_detalhados": []
        }
        
        self.setup_logging()
        self.criar_estrutura_diretorios()
    
    def setup_logging(self):
        """Configura o sistema de logging."""
        self.diretorio_logs.mkdir(parents=True, exist_ok=True)
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        log_file = self.diretorio_logs / f"consolidacao_{timestamp}.log"
        
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(levelname)s - %(message)s',
            handlers=[
                logging.FileHandler(log_file, encoding='utf-8'),
                logging.StreamHandler()
            ]
        )
        
        self.logger = logging.getLogger(__name__)
        self.logger.info("Iniciando processo de consolidação de arquivos")
    
    def criar_estrutura_diretorios(self):
        """Cria a estrutura de diretórios necessária."""
        self.diretorio_consolidado.mkdir(exist_ok=True)
        self.diretorio_txt.mkdir(exist_ok=True)
        self.logger.info(f"Estrutura de diretórios criada em: {self.diretorio_consolidado}")
    
    def calcular_checksum(self, arquivo: Path) -> str:
        """Calcula o checksum MD5 de um arquivo."""
        hash_md5 = hashlib.md5()
        try:
            with open(arquivo, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_md5.update(chunk)
            return hash_md5.hexdigest()
        except Exception as e:
            self.logger.error(f"Erro ao calcular checksum de {arquivo}: {e}")
            return ""
    
    def obter_metadados_arquivo(self, arquivo: Path) -> Dict:
        """Obtém metadados básicos do arquivo."""
        try:
            stat = arquivo.stat()
            return {
                "nome": arquivo.name,
                "tamanho": stat.st_size,
                "data_criacao": datetime.fromtimestamp(stat.st_ctime).isoformat(),
                "data_modificacao": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                "checksum": self.calcular_checksum(arquivo)
            }
        except Exception as e:
            self.logger.error(f"Erro ao obter metadados de {arquivo}: {e}")
            return {}
    
    def listar_todos_arquivos(self) -> List[Tuple[Path, str]]:
        """Lista todos os arquivos das pastas especificadas."""
        arquivos = []
        
        # Arquivo individual na raiz
        arquivo_raiz = self.diretorio_base / "_[FSS] SWIPE FILE COMERCIAL.docx"
        if arquivo_raiz.exists():
            arquivos.append((arquivo_raiz, "Arquivo Raiz"))
        
        # Arquivos das pastas categorizadas
        for pasta, categoria in self.categorias.items():
            pasta_path = self.diretorio_base / pasta
            if pasta_path.exists():
                for arquivo in pasta_path.rglob("*"):
                    if arquivo.is_file():
                        arquivos.append((arquivo, categoria))

        # Fallback: se nenhuma categoria esperada for encontrada, varre toda a árvore
        if not arquivos:
            for arquivo in self.diretorio_base.rglob("*"):
                if arquivo.is_file() and "Consolidado" not in str(arquivo) and not arquivo.name.startswith('.'):
                    arquivos.append((arquivo, "Outros"))

        self.estatisticas["total_arquivos"] = len(arquivos)
        self.logger.info(f"Total de arquivos encontrados: {len(arquivos)}")
        return arquivos
    
    def converter_docx_para_txt(self, arquivo: Path) -> str:
        """Converte arquivo DOCX para texto."""
        if not DOCX_AVAILABLE:
            raise ImportError("Biblioteca python-docx não disponível")
        
        try:
            doc = Document(arquivo)
            texto = []
            for paragrafo in doc.paragraphs:
                if paragrafo.text.strip():
                    texto.append(paragrafo.text)
            return "\n".join(texto)
        except Exception as e:
            raise Exception(f"Erro na conversão DOCX: {e}")
    
    def converter_pdf_para_txt(self, arquivo: Path) -> str:
        """Converte arquivo PDF para texto."""
        if not PDF_AVAILABLE:
            raise ImportError("Biblioteca PyPDF2 não disponível")
        
        try:
            texto = []
            with open(arquivo, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    texto.append(page.extract_text())
            return "\n".join(texto)
        except Exception as e:
            raise Exception(f"Erro na conversão PDF: {e}")
    
    def converter_xlsx_para_txt(self, arquivo: Path) -> str:
        """Converte arquivo XLSX para texto."""
        if not XLSX_AVAILABLE:
            raise ImportError("Biblioteca openpyxl não disponível")
        
        try:
            workbook = openpyxl.load_workbook(arquivo)
            texto = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                texto.append(f"=== PLANILHA: {sheet_name} ===\n")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        texto.append(row_text)
                texto.append("\n")
            
            return "\n".join(texto)
        except Exception as e:
            raise Exception(f"Erro na conversão XLSX: {e}")
    
    def converter_pptx_para_txt(self, arquivo: Path) -> str:
        """Converte arquivo PPTX para texto."""
        if not PPTX_AVAILABLE:
            raise ImportError("Biblioteca python-pptx não disponível")
        
        try:
            prs = Presentation(arquivo)
            texto = []
            
            for i, slide in enumerate(prs.slides, 1):
                texto.append(f"=== SLIDE {i} ===\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texto.append(shape.text)
                texto.append("\n")
            
            return "\n".join(texto)
        except Exception as e:
            raise Exception(f"Erro na conversão PPTX: {e}")
    
    def converter_arquivo_para_txt(self, arquivo: Path) -> Optional[str]:
        """Converte arquivo para texto baseado na extensão."""
        extensao = arquivo.suffix.lower()
        
        try:
            if extensao == '.docx':
                return self.converter_docx_para_txt(arquivo)
            elif extensao == '.pdf':
                return self.converter_pdf_para_txt(arquivo)
            elif extensao == '.xlsx':
                return self.converter_xlsx_para_txt(arquivo)
            elif extensao == '.pptx':
                return self.converter_pptx_para_txt(arquivo)
            elif extensao == '.txt':
                with open(arquivo, 'r', encoding='utf-8') as f:
                    return f.read()
            else:
                # Para outros formatos, tenta ler como texto
                try:
                    with open(arquivo, 'r', encoding='utf-8') as f:
                        return f.read()
                except UnicodeDecodeError:
                    # Se não conseguir ler como UTF-8, tenta outras codificações
                    for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                        try:
                            with open(arquivo, 'r', encoding=encoding) as f:
                                return f.read()
                        except UnicodeDecodeError:
                            continue
                    raise Exception("Não foi possível decodificar o arquivo")
        
        except Exception as e:
            self.logger.error(f"Erro na conversão de {arquivo}: {e}")
            self.estatisticas["conversoes_erro"] += 1
            self.estatisticas["erros_detalhados"].append({
                "arquivo": str(arquivo),
                "erro": str(e),
                "tipo": "conversao"
            })
            return None
    
    def processar_arquivos(self):
        """Processa todos os arquivos: copia, renomeia e converte."""
        arquivos = self.listar_todos_arquivos()
        contador = 1
        
        for arquivo_path, categoria in arquivos:
            try:
                # Atualiza estatísticas de formato
                extensao = arquivo_path.suffix.lower()
                self.estatisticas["formatos_encontrados"][extensao] = \
                    self.estatisticas["formatos_encontrados"].get(extensao, 0) + 1
                
                # Gera nome consolidado
                nome_original = arquivo_path.stem
                extensao_original = arquivo_path.suffix
                nome_consolidado = f"{contador:03d}. {nome_original} ({categoria}){extensao_original}"
                
                # Copia arquivo para pasta consolidada
                arquivo_consolidado = self.diretorio_consolidado / nome_consolidado
                shutil.copy2(arquivo_path, arquivo_consolidado)
                
                # Obtém metadados
                metadados = self.obter_metadados_arquivo(arquivo_path)
                
                # Converte para TXT
                conteudo_txt = self.converter_arquivo_para_txt(arquivo_path)
                if conteudo_txt is not None:
                    # Salva arquivo TXT
                    nome_txt = f"{contador:03d}. {nome_original} ({categoria}).txt"
                    arquivo_txt = self.diretorio_txt / nome_txt
                    
                    with open(arquivo_txt, 'w', encoding='utf-8') as f:
                        # Adiciona metadados no início do arquivo
                        f.write(f"=== METADADOS ===\n")
                        f.write(f"Arquivo Original: {arquivo_path}\n")
                        f.write(f"Categoria: {categoria}\n")
                        f.write(f"Data de Processamento: {datetime.now().isoformat()}\n")
                        for key, value in metadados.items():
                            f.write(f"{key.title()}: {value}\n")
                        f.write(f"\n=== CONTEÚDO ===\n\n")
                        f.write(conteudo_txt)
                    
                    self.estatisticas["conversoes_sucesso"] += 1
                    self.logger.info(f"Processado com sucesso: {nome_consolidado}")
                else:
                    self.logger.warning(f"Falha na conversão: {nome_consolidado}")
                
                self.estatisticas["arquivos_processados"] += 1
                contador += 1
                
            except Exception as e:
                self.estatisticas["arquivos_com_erro"] += 1
                self.estatisticas["erros_detalhados"].append({
                    "arquivo": str(arquivo_path),
                    "erro": str(e),
                    "traceback": traceback.format_exc(),
                    "tipo": "processamento"
                })
                self.logger.error(f"Erro ao processar {arquivo_path}: {e}")
                contador += 1
    
    def gerar_relatorio_validacao(self):
        """Gera relatório final de validação."""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        relatorio_path = self.diretorio_logs / f"relatorio_validacao_{timestamp}.json"
        
        # Calcula taxa de sucesso
        if self.estatisticas["total_arquivos"] > 0:
            taxa_sucesso = (self.estatisticas["arquivos_processados"] / 
                          self.estatisticas["total_arquivos"]) * 100
            taxa_conversao = (self.estatisticas["conversoes_sucesso"] / 
                            self.estatisticas["total_arquivos"]) * 100
        else:
            taxa_sucesso = 0
            taxa_conversao = 0
        
        relatorio = {
            "timestamp": datetime.now().isoformat(),
            "estatisticas": self.estatisticas,
            "taxa_sucesso_processamento": round(taxa_sucesso, 2),
            "taxa_sucesso_conversao": round(taxa_conversao, 2),
            "bibliotecas_disponiveis": {
                "docx": DOCX_AVAILABLE,
                "pdf": PDF_AVAILABLE,
                "xlsx": XLSX_AVAILABLE,
                "pptx": PPTX_AVAILABLE
            }
        }
        
        with open(relatorio_path, 'w', encoding='utf-8') as f:
            json.dump(relatorio, f, indent=2, ensure_ascii=False)
        
        # Log do resumo
        self.logger.info("=== RELATÓRIO FINAL ===")
        self.logger.info(f"Total de arquivos: {self.estatisticas['total_arquivos']}")
        self.logger.info(f"Arquivos processados: {self.estatisticas['arquivos_processados']}")
        self.logger.info(f"Arquivos com erro: {self.estatisticas['arquivos_com_erro']}")
        self.logger.info(f"Conversões bem-sucedidas: {self.estatisticas['conversoes_sucesso']}")
        self.logger.info(f"Conversões com erro: {self.estatisticas['conversoes_erro']}")
        self.logger.info(f"Taxa de sucesso: {taxa_sucesso:.2f}%")
        self.logger.info(f"Taxa de conversão: {taxa_conversao:.2f}%")
        self.logger.info(f"Relatório salvo em: {relatorio_path}")
        
        return relatorio
    
    def executar(self):
        """Executa o processo completo de consolidação."""
        try:
            self.logger.info("Iniciando processo de consolidação...")
            self.processar_arquivos()
            relatorio = self.gerar_relatorio_validacao()
            self.logger.info("Processo de consolidação concluído com sucesso!")
            return relatorio
        except Exception as e:
            self.logger.error(f"Erro crítico no processo: {e}")
            self.logger.error(traceback.format_exc())
            raise


def main():
    """Função principal."""
    import argparse
    parser = argparse.ArgumentParser(description="Consolidador de Arquivos")
    parser.add_argument(
        "--diretorio", "-d",
        default="/Users/marcosdaniels/Downloads/5. FULL SALES 2",
        help="Diretório base para consolidação"
    )
    args = parser.parse_args()
    diretorio_base = args.diretorio

    print("=== CONSOLIDADOR DE ARQUIVOS ===")
    print(f"Diretório base: {diretorio_base}")
    print("\nVerificando bibliotecas disponíveis...")
    print(f"DOCX: {'✓' if DOCX_AVAILABLE else '✗'}")
    print(f"PDF: {'✓' if PDF_AVAILABLE else '✗'}")
    print(f"XLSX: {'✓' if XLSX_AVAILABLE else '✗'}")
    print(f"PPTX: {'✓' if PPTX_AVAILABLE else '✗'}")
    
    if not any([DOCX_AVAILABLE, PDF_AVAILABLE, XLSX_AVAILABLE, PPTX_AVAILABLE]):
        print("\n⚠️  AVISO: Nenhuma biblioteca de conversão está disponível.")
        print("Instale as dependências necessárias:")
        print("pip install python-docx PyPDF2 openpyxl python-pptx")
        print("\nContinuando apenas com cópia de arquivos...")
    
    consolidador = ConsolidadorArquivos(diretorio_base)
    relatorio = consolidador.executar()
    
    print(f"\n=== PROCESSO CONCLUÍDO ===")
    print(f"Arquivos processados: {relatorio['estatisticas']['arquivos_processados']}")
    print(f"Taxa de sucesso: {relatorio['taxa_sucesso_processamento']:.2f}%")
    print(f"Arquivos TXT gerados: {relatorio['estatisticas']['conversoes_sucesso']}")


if __name__ == "__main__":
    main()