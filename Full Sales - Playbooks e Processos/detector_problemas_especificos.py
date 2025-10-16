#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Detector de Problemas Específicos em Arquivos
Detecta problemas sutis que podem não ser identificados por verificações básicas
Autor: Sistema de Consolidação de Arquivos
Data: 2025-10-02
"""

import os
import sys
import json
import logging
from datetime import datetime
from pathlib import Path
import hashlib
import tempfile
from typing import Dict, List, Tuple, Optional
import zipfile
import xml.etree.ElementTree as ET

# Importações para diferentes tipos de arquivo
try:
    from docx import Document
    from docx.opc.exceptions import PackageNotFoundError
    import docx2txt
except ImportError:
    print("Instalando python-docx e docx2txt...")
    os.system("pip install python-docx docx2txt")
    from docx import Document
    from docx.opc.exceptions import PackageNotFoundError
    import docx2txt

try:
    import openpyxl
    from openpyxl.utils.exceptions import InvalidFileException
except ImportError:
    print("Instalando openpyxl...")
    os.system("pip install openpyxl")
    import openpyxl
    from openpyxl.utils.exceptions import InvalidFileException

try:
    import PyPDF2
    from PyPDF2.errors import PdfReadError
except ImportError:
    print("Instalando PyPDF2...")
    os.system("pip install PyPDF2")
    import PyPDF2
    from PyPDF2.errors import PdfReadError

try:
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError as PPTXPackageNotFoundError
except ImportError:
    print("Instalando python-pptx...")
    os.system("pip install python-pptx")
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError as PPTXPackageNotFoundError

class DetectorProblemasEspecificos:
    """Classe para detectar problemas específicos e sutis em arquivos."""
    
    def __init__(self, diretorio_base: str):
        self.diretorio_base = diretorio_base
        self.diretorio_logs = os.path.join(diretorio_base, "logs")
        os.makedirs(self.diretorio_logs, exist_ok=True)
        
        # Configura logger
        self.logger = self._setup_logger()
        
        self.problemas_encontrados = {
            'arquivos_vazios': [],
            'arquivos_muito_pequenos': [],
            'arquivos_com_conteudo_suspeito': [],
            'arquivos_com_estrutura_danificada': [],
            'arquivos_com_metadados_corrompidos': [],
            'arquivos_com_encoding_problemas': [],
            'arquivos_duplicados': [],
            'arquivos_com_extensao_incorreta': []
        }
    
    def _setup_logger(self) -> logging.Logger:
        """Configura o sistema de logging."""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        log_file = os.path.join(self.diretorio_logs, f"detector_problemas_{timestamp}.log")
        
        logger = logging.getLogger('DetectorProblemas')
        logger.setLevel(logging.INFO)
        
        # Remove handlers existentes
        for handler in logger.handlers[:]:
            logger.removeHandler(handler)
        
        # Handler para arquivo
        file_handler = logging.FileHandler(log_file, encoding='utf-8')
        file_handler.setLevel(logging.INFO)
        
        # Handler para console
        console_handler = logging.StreamHandler()
        console_handler.setLevel(logging.INFO)
        
        # Formato
        formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
        file_handler.setFormatter(formatter)
        console_handler.setFormatter(formatter)
        
        logger.addHandler(file_handler)
        logger.addHandler(console_handler)
        
        return logger
    
    def executar_deteccao_completa(self) -> Dict:
        """Executa detecção completa de problemas específicos."""
        self.logger.info("Iniciando detecção de problemas específicos...")
        
        # Encontra todos os arquivos
        arquivos = self._encontrar_todos_arquivos()
        
        resultados = {
            'timestamp': datetime.now().strftime("%Y%m%d_%H%M%S"),
            'total_arquivos': len(arquivos),
            'arquivos_analisados': 0,
            'problemas_encontrados': {},
            'detalhes_problemas': [],
            'recomendacoes': [],
            'estatisticas': {}
        }
        
        # Executa diferentes tipos de verificação
        self._verificar_arquivos_vazios(arquivos)
        self._verificar_arquivos_pequenos(arquivos)
        self._verificar_conteudo_suspeito(arquivos)
        self._verificar_estrutura_interna(arquivos)
        self._verificar_metadados(arquivos)
        self._verificar_encoding(arquivos)
        self._verificar_duplicados(arquivos)
        self._verificar_extensoes_incorretas(arquivos)
        
        # Compila resultados
        resultados['problemas_encontrados'] = self.problemas_encontrados
        resultados['arquivos_analisados'] = len(arquivos)
        
        # Gera estatísticas e recomendações
        resultados['estatisticas'] = self._gerar_estatisticas()
        resultados['recomendacoes'] = self._gerar_recomendacoes()
        
        # Salva relatório
        self._salvar_relatorio(resultados)
        
        # Exibe resumo
        self._exibir_resumo(resultados)
        
        return resultados
    
    def _encontrar_todos_arquivos(self) -> List[str]:
        """Encontra todos os arquivos para análise."""
        arquivos = []
        
        for root, dirs, files in os.walk(self.diretorio_base):
            # Pula diretórios de sistema
            if any(skip in root for skip in ['logs', 'backups', '__pycache__']):
                continue
                
            for file in files:
                arquivo_path = os.path.join(root, file)
                arquivos.append(arquivo_path)
        
        return arquivos
    
    def _verificar_arquivos_vazios(self, arquivos: List[str]):
        """Verifica arquivos completamente vazios."""
        self.logger.info("Verificando arquivos vazios...")
        
        for arquivo in arquivos:
            try:
                if os.path.getsize(arquivo) == 0:
                    self.problemas_encontrados['arquivos_vazios'].append({
                        'arquivo': arquivo,
                        'problema': 'Arquivo completamente vazio',
                        'tamanho': 0
                    })
                    self.logger.warning(f"Arquivo vazio encontrado: {os.path.basename(arquivo)}")
            except Exception as e:
                self.logger.error(f"Erro ao verificar tamanho de {arquivo}: {str(e)}")
    
    def _verificar_arquivos_pequenos(self, arquivos: List[str]):
        """Verifica arquivos suspeitosamente pequenos para seu tipo."""
        self.logger.info("Verificando arquivos suspeitosamente pequenos...")
        
        tamanhos_minimos = {
            '.docx': 5000,    # 5KB mínimo para DOCX
            '.xlsx': 3000,    # 3KB mínimo para XLSX
            '.pptx': 10000,   # 10KB mínimo para PPTX
            '.pdf': 1000      # 1KB mínimo para PDF
        }
        
        for arquivo in arquivos:
            try:
                extensao = Path(arquivo).suffix.lower()
                tamanho = os.path.getsize(arquivo)
                
                if extensao in tamanhos_minimos and tamanho < tamanhos_minimos[extensao]:
                    self.problemas_encontrados['arquivos_muito_pequenos'].append({
                        'arquivo': arquivo,
                        'problema': f'Arquivo {extensao} muito pequeno',
                        'tamanho': tamanho,
                        'tamanho_minimo_esperado': tamanhos_minimos[extensao]
                    })
                    self.logger.warning(f"Arquivo pequeno: {os.path.basename(arquivo)} ({tamanho} bytes)")
            except Exception as e:
                self.logger.error(f"Erro ao verificar tamanho de {arquivo}: {str(e)}")
    
    def _verificar_conteudo_suspeito(self, arquivos: List[str]):
        """Verifica conteúdo suspeito nos arquivos."""
        self.logger.info("Verificando conteúdo suspeito...")
        
        for arquivo in arquivos:
            try:
                extensao = Path(arquivo).suffix.lower()
                
                if extensao == '.docx':
                    self._verificar_conteudo_docx(arquivo)
                elif extensao == '.xlsx':
                    self._verificar_conteudo_xlsx(arquivo)
                elif extensao == '.pptx':
                    self._verificar_conteudo_pptx(arquivo)
                elif extensao == '.pdf':
                    self._verificar_conteudo_pdf(arquivo)
                    
            except Exception as e:
                self.logger.error(f"Erro ao verificar conteúdo de {arquivo}: {str(e)}")
    
    def _verificar_conteudo_docx(self, arquivo: str):
        """Verifica conteúdo específico de DOCX."""
        try:
            doc = Document(arquivo)
            
            # Verifica se tem parágrafos
            if len(doc.paragraphs) == 0:
                self.problemas_encontrados['arquivos_com_conteudo_suspeito'].append({
                    'arquivo': arquivo,
                    'problema': 'DOCX sem parágrafos',
                    'detalhes': 'Documento não contém texto'
                })
            
            # Verifica se todos os parágrafos estão vazios
            texto_total = ''.join([p.text for p in doc.paragraphs]).strip()
            if not texto_total:
                self.problemas_encontrados['arquivos_com_conteudo_suspeito'].append({
                    'arquivo': arquivo,
                    'problema': 'DOCX com parágrafos vazios',
                    'detalhes': 'Todos os parágrafos estão vazios'
                })
                
        except Exception as e:
            self.problemas_encontrados['arquivos_com_conteudo_suspeito'].append({
                'arquivo': arquivo,
                'problema': 'Erro ao ler DOCX',
                'detalhes': str(e)
            })
    
    def _verificar_conteudo_xlsx(self, arquivo: str):
        """Verifica conteúdo específico de XLSX."""
        try:
            wb = openpyxl.load_workbook(arquivo, read_only=True)
            
            # Verifica se tem planilhas
            if len(wb.sheetnames) == 0:
                self.problemas_encontrados['arquivos_com_conteudo_suspeito'].append({
                    'arquivo': arquivo,
                    'problema': 'XLSX sem planilhas',
                    'detalhes': 'Arquivo não contém planilhas'
                })
            
            # Verifica se as planilhas têm dados
            planilhas_vazias = 0
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                if ws.max_row <= 1 and ws.max_column <= 1:
                    planilhas_vazias += 1
            
            if planilhas_vazias == len(wb.sheetnames):
                self.problemas_encontrados['arquivos_com_conteudo_suspeito'].append({
                    'arquivo': arquivo,
                    'problema': 'XLSX com planilhas vazias',
                    'detalhes': 'Todas as planilhas estão vazias'
                })
            
            wb.close()
            
        except Exception as e:
            self.problemas_encontrados['arquivos_com_conteudo_suspeito'].append({
                'arquivo': arquivo,
                'problema': 'Erro ao ler XLSX',
                'detalhes': str(e)
            })
    
    def _verificar_conteudo_pptx(self, arquivo: str):
        """Verifica conteúdo específico de PPTX."""
        try:
            prs = Presentation(arquivo)
            
            # Verifica se tem slides
            if len(prs.slides) == 0:
                self.problemas_encontrados['arquivos_com_conteudo_suspeito'].append({
                    'arquivo': arquivo,
                    'problema': 'PPTX sem slides',
                    'detalhes': 'Apresentação não contém slides'
                })
            
            # Verifica se os slides têm conteúdo
            slides_vazios = 0
            for slide in prs.slides:
                tem_conteudo = False
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        tem_conteudo = True
                        break
                if not tem_conteudo:
                    slides_vazios += 1
            
            if slides_vazios == len(prs.slides):
                self.problemas_encontrados['arquivos_com_conteudo_suspeito'].append({
                    'arquivo': arquivo,
                    'problema': 'PPTX com slides vazios',
                    'detalhes': 'Todos os slides estão vazios'
                })
                
        except Exception as e:
            self.problemas_encontrados['arquivos_com_conteudo_suspeito'].append({
                'arquivo': arquivo,
                'problema': 'Erro ao ler PPTX',
                'detalhes': str(e)
            })
    
    def _verificar_conteudo_pdf(self, arquivo: str):
        """Verifica conteúdo específico de PDF."""
        try:
            with open(arquivo, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                
                # Verifica se tem páginas
                if len(reader.pages) == 0:
                    self.problemas_encontrados['arquivos_com_conteudo_suspeito'].append({
                        'arquivo': arquivo,
                        'problema': 'PDF sem páginas',
                        'detalhes': 'PDF não contém páginas'
                    })
                
                # Verifica se as páginas têm texto
                paginas_sem_texto = 0
                for page in reader.pages:
                    try:
                        texto = page.extract_text().strip()
                        if not texto:
                            paginas_sem_texto += 1
                    except:
                        paginas_sem_texto += 1
                
                if paginas_sem_texto == len(reader.pages) and len(reader.pages) > 0:
                    self.problemas_encontrados['arquivos_com_conteudo_suspeito'].append({
                        'arquivo': arquivo,
                        'problema': 'PDF sem texto extraível',
                        'detalhes': 'Todas as páginas estão sem texto ou são imagens'
                    })
                    
        except Exception as e:
            self.problemas_encontrados['arquivos_com_conteudo_suspeito'].append({
                'arquivo': arquivo,
                'problema': 'Erro ao ler PDF',
                'detalhes': str(e)
            })
    
    def _verificar_estrutura_interna(self, arquivos: List[str]):
        """Verifica estrutura interna dos arquivos Office."""
        self.logger.info("Verificando estrutura interna...")
        
        for arquivo in arquivos:
            try:
                extensao = Path(arquivo).suffix.lower()
                
                if extensao in ['.docx', '.xlsx', '.pptx']:
                    self._verificar_estrutura_zip(arquivo)
                    
            except Exception as e:
                self.logger.error(f"Erro ao verificar estrutura de {arquivo}: {str(e)}")
    
    def _verificar_estrutura_zip(self, arquivo: str):
        """Verifica estrutura ZIP dos arquivos Office."""
        try:
            with zipfile.ZipFile(arquivo, 'r') as zip_file:
                arquivos_internos = zip_file.namelist()
                
                # Verifica arquivos essenciais
                arquivos_essenciais = {
                    '.docx': ['word/document.xml', '[Content_Types].xml'],
                    '.xlsx': ['xl/workbook.xml', '[Content_Types].xml'],
                    '.pptx': ['ppt/presentation.xml', '[Content_Types].xml']
                }
                
                extensao = Path(arquivo).suffix.lower()
                if extensao in arquivos_essenciais:
                    for essencial in arquivos_essenciais[extensao]:
                        if essencial not in arquivos_internos:
                            self.problemas_encontrados['arquivos_com_estrutura_danificada'].append({
                                'arquivo': arquivo,
                                'problema': f'Arquivo essencial ausente: {essencial}',
                                'detalhes': f'Estrutura {extensao} incompleta'
                            })
                
                # Verifica se consegue ler os XMLs principais
                for arquivo_interno in arquivos_internos:
                    if arquivo_interno.endswith('.xml'):
                        try:
                            conteudo = zip_file.read(arquivo_interno)
                            ET.fromstring(conteudo)
                        except ET.ParseError:
                            self.problemas_encontrados['arquivos_com_estrutura_danificada'].append({
                                'arquivo': arquivo,
                                'problema': f'XML corrompido: {arquivo_interno}',
                                'detalhes': 'Estrutura XML inválida'
                            })
                        except Exception:
                            pass  # Ignora outros erros de leitura
                            
        except zipfile.BadZipFile:
            self.problemas_encontrados['arquivos_com_estrutura_danificada'].append({
                'arquivo': arquivo,
                'problema': 'Estrutura ZIP corrompida',
                'detalhes': 'Arquivo não é um ZIP válido'
            })
        except Exception as e:
            self.problemas_encontrados['arquivos_com_estrutura_danificada'].append({
                'arquivo': arquivo,
                'problema': 'Erro na verificação de estrutura',
                'detalhes': str(e)
            })
    
    def _verificar_metadados(self, arquivos: List[str]):
        """Verifica problemas nos metadados."""
        self.logger.info("Verificando metadados...")
        
        for arquivo in arquivos:
            try:
                # Verifica se consegue obter informações básicas do arquivo
                stat = os.stat(arquivo)
                
                # Verifica datas suspeitas
                if stat.st_mtime < 0 or stat.st_ctime < 0:
                    self.problemas_encontrados['arquivos_com_metadados_corrompidos'].append({
                        'arquivo': arquivo,
                        'problema': 'Datas de arquivo inválidas',
                        'detalhes': 'Timestamps negativos ou inválidos'
                    })
                
                # Verifica permissões
                if not os.access(arquivo, os.R_OK):
                    self.problemas_encontrados['arquivos_com_metadados_corrompidos'].append({
                        'arquivo': arquivo,
                        'problema': 'Arquivo não legível',
                        'detalhes': 'Sem permissão de leitura'
                    })
                    
            except Exception as e:
                self.problemas_encontrados['arquivos_com_metadados_corrompidos'].append({
                    'arquivo': arquivo,
                    'problema': 'Erro ao acessar metadados',
                    'detalhes': str(e)
                })
    
    def _verificar_encoding(self, arquivos: List[str]):
        """Verifica problemas de encoding."""
        self.logger.info("Verificando problemas de encoding...")
        
        for arquivo in arquivos:
            try:
                # Verifica se o nome do arquivo tem caracteres problemáticos
                nome_arquivo = os.path.basename(arquivo)
                
                # Tenta codificar/decodificar o nome
                try:
                    nome_arquivo.encode('utf-8').decode('utf-8')
                except UnicodeError:
                    self.problemas_encontrados['arquivos_com_encoding_problemas'].append({
                        'arquivo': arquivo,
                        'problema': 'Nome de arquivo com encoding inválido',
                        'detalhes': 'Caracteres não UTF-8 no nome'
                    })
                
                # Verifica caracteres de controle no nome
                if any(ord(c) < 32 for c in nome_arquivo if c != '\n' and c != '\r' and c != '\t'):
                    self.problemas_encontrados['arquivos_com_encoding_problemas'].append({
                        'arquivo': arquivo,
                        'problema': 'Caracteres de controle no nome',
                        'detalhes': 'Nome contém caracteres não imprimíveis'
                    })
                    
            except Exception as e:
                self.problemas_encontrados['arquivos_com_encoding_problemas'].append({
                    'arquivo': arquivo,
                    'problema': 'Erro na verificação de encoding',
                    'detalhes': str(e)
                })
    
    def _verificar_duplicados(self, arquivos: List[str]):
        """Verifica arquivos duplicados por checksum."""
        self.logger.info("Verificando arquivos duplicados...")
        
        checksums = {}
        
        for arquivo in arquivos:
            try:
                checksum = self._calcular_checksum(arquivo)
                if checksum:
                    if checksum in checksums:
                        self.problemas_encontrados['arquivos_duplicados'].append({
                            'arquivo': arquivo,
                            'problema': 'Arquivo duplicado',
                            'detalhes': f'Idêntico a: {checksums[checksum]}',
                            'checksum': checksum
                        })
                    else:
                        checksums[checksum] = arquivo
                        
            except Exception as e:
                self.logger.error(f"Erro ao calcular checksum de {arquivo}: {str(e)}")
    
    def _verificar_extensoes_incorretas(self, arquivos: List[str]):
        """Verifica arquivos com extensões incorretas."""
        self.logger.info("Verificando extensões incorretas...")
        
        for arquivo in arquivos:
            try:
                extensao = Path(arquivo).suffix.lower()
                
                # Lê os primeiros bytes para identificar o tipo real
                with open(arquivo, 'rb') as f:
                    header = f.read(16)
                
                tipo_real = self._identificar_tipo_por_header(header)
                
                if tipo_real and tipo_real != extensao:
                    self.problemas_encontrados['arquivos_com_extensao_incorreta'].append({
                        'arquivo': arquivo,
                        'problema': 'Extensão incorreta',
                        'detalhes': f'Extensão: {extensao}, Tipo real: {tipo_real}',
                        'extensao_atual': extensao,
                        'tipo_real': tipo_real
                    })
                    
            except Exception as e:
                self.logger.error(f"Erro ao verificar extensão de {arquivo}: {str(e)}")
    
    def _identificar_tipo_por_header(self, header: bytes) -> Optional[str]:
        """Identifica tipo de arquivo pelo header."""
        if header.startswith(b'PK\x03\x04'):
            # Arquivo ZIP (pode ser DOCX, XLSX, PPTX)
            return '.zip_based'
        elif header.startswith(b'%PDF-'):
            return '.pdf'
        elif header.startswith(b'\xff\xd8\xff'):
            return '.jpg'
        elif header.startswith(b'\x89PNG'):
            return '.png'
        elif header.startswith(b'GIF8'):
            return '.gif'
        
        return None
    
    def _calcular_checksum(self, arquivo: str) -> Optional[str]:
        """Calcula checksum MD5 do arquivo."""
        try:
            hash_md5 = hashlib.md5()
            with open(arquivo, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_md5.update(chunk)
            return hash_md5.hexdigest()
        except:
            return None
    
    def _gerar_estatisticas(self) -> Dict:
        """Gera estatísticas dos problemas encontrados."""
        estatisticas = {}
        
        for tipo_problema, problemas in self.problemas_encontrados.items():
            estatisticas[tipo_problema] = {
                'total': len(problemas),
                'arquivos': [p['arquivo'] for p in problemas]
            }
        
        return estatisticas
    
    def _gerar_recomendacoes(self) -> List[str]:
        """Gera recomendações baseadas nos problemas encontrados."""
        recomendacoes = []
        
        if self.problemas_encontrados['arquivos_vazios']:
            recomendacoes.append("Remover ou investigar arquivos completamente vazios")
        
        if self.problemas_encontrados['arquivos_muito_pequenos']:
            recomendacoes.append("Verificar arquivos suspeitosamente pequenos - podem estar corrompidos")
        
        if self.problemas_encontrados['arquivos_com_conteudo_suspeito']:
            recomendacoes.append("Revisar arquivos com conteúdo suspeito ou vazio")
        
        if self.problemas_encontrados['arquivos_com_estrutura_danificada']:
            recomendacoes.append("Tentar reparar arquivos com estrutura danificada")
        
        if self.problemas_encontrados['arquivos_duplicados']:
            recomendacoes.append("Considerar remoção de arquivos duplicados para economizar espaço")
        
        if self.problemas_encontrados['arquivos_com_extensao_incorreta']:
            recomendacoes.append("Corrigir extensões de arquivos incorretas")
        
        return recomendacoes
    
    def _salvar_relatorio(self, resultados: Dict):
        """Salva relatório detalhado."""
        timestamp = resultados['timestamp']
        arquivo_relatorio = os.path.join(self.diretorio_logs, f"relatorio_problemas_especificos_{timestamp}.json")
        
        try:
            with open(arquivo_relatorio, 'w', encoding='utf-8') as f:
                json.dump(resultados, f, indent=2, ensure_ascii=False)
            
            self.logger.info(f"Relatório salvo em: {arquivo_relatorio}")
            
        except Exception as e:
            self.logger.error(f"Erro ao salvar relatório: {str(e)}")
    
    def _exibir_resumo(self, resultados: Dict):
        """Exibe resumo dos resultados."""
        print("\n" + "="*60)
        print("RELATÓRIO - DETECÇÃO DE PROBLEMAS ESPECÍFICOS")
        print("="*60)
        print(f"Total de arquivos analisados: {resultados['arquivos_analisados']}")
        
        total_problemas = sum(len(problemas) for problemas in self.problemas_encontrados.values())
        print(f"Total de problemas encontrados: {total_problemas}")
        
        for tipo_problema, problemas in self.problemas_encontrados.items():
            if problemas:
                print(f"\n{tipo_problema.replace('_', ' ').title()}: {len(problemas)}")
                for problema in problemas[:3]:  # Mostra apenas os primeiros 3
                    print(f"  - {os.path.basename(problema['arquivo'])}: {problema['problema']}")
                if len(problemas) > 3:
                    print(f"  ... e mais {len(problemas) - 3} arquivo(s)")
        
        if resultados['recomendacoes']:
            print("\nRecomendações:")
            for i, recomendacao in enumerate(resultados['recomendacoes'], 1):
                print(f"{i}. {recomendacao}")
        
        print("="*60)

def main():
    """Função principal."""
    # Diretório base
    diretorio_base = "/Users/marcosdaniels/Downloads/5. FULL SALES 2/Consolidado"
    
    if not os.path.exists(diretorio_base):
        print(f"Erro: Diretório {diretorio_base} não encontrado!")
        return
    
    # Cria e executa o detector
    detector = DetectorProblemasEspecificos(diretorio_base)
    resultados = detector.executar_deteccao_completa()
    
    print(f"\nProcesso concluído! Verifique os logs em: {detector.diretorio_logs}")

if __name__ == "__main__":
    main()