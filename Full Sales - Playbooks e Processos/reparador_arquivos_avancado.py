#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script Avançado para Detecção e Reparo de Arquivos Corrompidos
Autor: Sistema de Consolidação de Arquivos
Data: 2025-10-02
"""

import os
import sys
import json
import shutil
import zipfile
import logging
from datetime import datetime
from pathlib import Path
import hashlib
import tempfile
from typing import Dict, List, Tuple, Optional

# Importações para diferentes tipos de arquivo
try:
    from docx import Document
    from docx.opc.exceptions import PackageNotFoundError
    import docx2txt
except ImportError:
    print("Instalando python-docx e docx2txt...")
    os.system("pip install python-docx docx2txt")
    from docx import Document
    from docx.opc.exceptions import PackageNotFoundError
    import docx2txt

try:
    import openpyxl
    from openpyxl.utils.exceptions import InvalidFileException
except ImportError:
    print("Instalando openpyxl...")
    os.system("pip install openpyxl")
    import openpyxl
    from openpyxl.utils.exceptions import InvalidFileException

try:
    import PyPDF2
    from PyPDF2.errors import PdfReadError
except ImportError:
    print("Instalando PyPDF2...")
    os.system("pip install PyPDF2")
    import PyPDF2
    from PyPDF2.errors import PdfReadError

try:
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError as PPTXPackageNotFoundError
except ImportError:
    print("Instalando python-pptx...")
    os.system("pip install python-pptx")
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError as PPTXPackageNotFoundError

try:
    import pandas as pd
except ImportError:
    print("Instalando pandas...")
    os.system("pip install pandas")
    import pandas as pd

class LoggerConfig:
    """Configuração de logging para o sistema."""
    
    @staticmethod
    def setup_logger(log_dir: str) -> logging.Logger:
        """Configura o sistema de logging."""
        os.makedirs(log_dir, exist_ok=True)
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        log_file = os.path.join(log_dir, f"reparacao_avancada_{timestamp}.log")
        
        logger = logging.getLogger('ReparadorAvancado')
        logger.setLevel(logging.INFO)
        
        # Remove handlers existentes
        for handler in logger.handlers[:]:
            logger.removeHandler(handler)
        
        # Handler para arquivo
        file_handler = logging.FileHandler(log_file, encoding='utf-8')
        file_handler.setLevel(logging.INFO)
        
        # Handler para console
        console_handler = logging.StreamHandler()
        console_handler.setLevel(logging.INFO)
        
        # Formato
        formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
        file_handler.setFormatter(formatter)
        console_handler.setFormatter(formatter)
        
        logger.addHandler(file_handler)
        logger.addHandler(console_handler)
        
        return logger

class DetectorCorrupcao:
    """Classe para detectar diferentes tipos de corrupção em arquivos."""
    
    def __init__(self, logger: logging.Logger):
        self.logger = logger
        self.resultados = {
            'arquivos_corrompidos': [],
            'arquivos_suspeitos': [],
            'arquivos_ok': [],
            'erros_deteccao': []
        }
    
    def verificar_integridade_basica(self, arquivo: str) -> Dict:
        """Verifica integridade básica do arquivo."""
        resultado = {
            'arquivo': arquivo,
            'tamanho': 0,
            'legivel': False,
            'extensao_valida': False,
            'checksum': None,
            'problemas': []
        }
        
        try:
            # Verifica se o arquivo existe e é legível
            if not os.path.exists(arquivo):
                resultado['problemas'].append('Arquivo não encontrado')
                return resultado
            
            # Verifica tamanho
            resultado['tamanho'] = os.path.getsize(arquivo)
            if resultado['tamanho'] == 0:
                resultado['problemas'].append('Arquivo vazio')
            
            # Verifica se é legível
            try:
                with open(arquivo, 'rb') as f:
                    f.read(1024)  # Tenta ler os primeiros 1KB
                resultado['legivel'] = True
            except Exception as e:
                resultado['problemas'].append(f'Erro de leitura: {str(e)}')
            
            # Verifica extensão
            extensao = Path(arquivo).suffix.lower()
            extensoes_validas = ['.docx', '.xlsx', '.pptx', '.pdf', '.txt', '.csv', '.png', '.jpg', '.jpeg']
            resultado['extensao_valida'] = extensao in extensoes_validas
            
            # Calcula checksum
            if resultado['legivel']:
                resultado['checksum'] = self._calcular_checksum(arquivo)
            
        except Exception as e:
            resultado['problemas'].append(f'Erro na verificação básica: {str(e)}')
            self.logger.error(f"Erro na verificação básica de {arquivo}: {str(e)}")
        
        return resultado
    
    def verificar_docx(self, arquivo: str) -> Dict:
        """Verifica integridade específica de arquivos DOCX."""
        resultado = {'tipo': 'DOCX', 'corrompido': False, 'problemas': [], 'reparavel': False}
        
        try:
            # Tenta abrir como documento Word
            doc = Document(arquivo)
            paragrafos = len(doc.paragraphs)
            resultado['paragrafos'] = paragrafos
            self.logger.info(f"DOCX OK: {os.path.basename(arquivo)} ({paragrafos} parágrafos)")
            
        except PackageNotFoundError:
            resultado['corrompido'] = True
            resultado['problemas'].append('Arquivo DOCX corrompido - estrutura ZIP inválida')
            resultado['reparavel'] = self._verificar_reparabilidade_zip(arquivo)
            
        except Exception as e:
            resultado['corrompido'] = True
            resultado['problemas'].append(f'Erro DOCX: {str(e)}')
            resultado['reparavel'] = self._verificar_reparabilidade_zip(arquivo)
        
        return resultado
    
    def verificar_xlsx(self, arquivo: str) -> Dict:
        """Verifica integridade específica de arquivos XLSX."""
        resultado = {'tipo': 'XLSX', 'corrompido': False, 'problemas': [], 'reparavel': False}
        
        try:
            # Tenta abrir como planilha Excel
            wb = openpyxl.load_workbook(arquivo, read_only=True)
            sheets = len(wb.sheetnames)
            resultado['planilhas'] = sheets
            wb.close()
            self.logger.info(f"XLSX OK: {os.path.basename(arquivo)} ({sheets} planilhas)")
            
        except InvalidFileException:
            resultado['corrompido'] = True
            resultado['problemas'].append('Arquivo XLSX corrompido - formato inválido')
            resultado['reparavel'] = self._verificar_reparabilidade_zip(arquivo)
            
        except Exception as e:
            resultado['corrompido'] = True
            resultado['problemas'].append(f'Erro XLSX: {str(e)}')
            resultado['reparavel'] = self._verificar_reparabilidade_zip(arquivo)
        
        return resultado
    
    def verificar_pptx(self, arquivo: str) -> Dict:
        """Verifica integridade específica de arquivos PPTX."""
        resultado = {'tipo': 'PPTX', 'corrompido': False, 'problemas': [], 'reparavel': False}
        
        try:
            # Tenta abrir como apresentação PowerPoint
            prs = Presentation(arquivo)
            slides = len(prs.slides)
            resultado['slides'] = slides
            self.logger.info(f"PPTX OK: {os.path.basename(arquivo)} ({slides} slides)")
            
        except PPTXPackageNotFoundError:
            resultado['corrompido'] = True
            resultado['problemas'].append('Arquivo PPTX corrompido - estrutura ZIP inválida')
            resultado['reparavel'] = self._verificar_reparabilidade_zip(arquivo)
            
        except Exception as e:
            resultado['corrompido'] = True
            resultado['problemas'].append(f'Erro PPTX: {str(e)}')
            resultado['reparavel'] = self._verificar_reparabilidade_zip(arquivo)
        
        return resultado
    
    def verificar_pdf(self, arquivo: str) -> Dict:
        """Verifica integridade específica de arquivos PDF."""
        resultado = {'tipo': 'PDF', 'corrompido': False, 'problemas': [], 'reparavel': False}
        
        try:
            # Tenta abrir como PDF
            with open(arquivo, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                paginas = len(reader.pages)
                resultado['paginas'] = paginas
                
                # Tenta ler a primeira página para verificar integridade
                if paginas > 0:
                    try:
                        primeira_pagina = reader.pages[0]
                        texto = primeira_pagina.extract_text()
                    except (AttributeError, Exception):
                        # Fallback para versões diferentes do PyPDF2
                        try:
                            texto = primeira_pagina.extractText()
                        except:
                            texto = ""  # Se não conseguir extrair texto, continua sem erro
                    
            self.logger.info(f"PDF OK: {os.path.basename(arquivo)} ({paginas} páginas)")
            
        except PdfReadError as e:
            resultado['corrompido'] = True
            resultado['problemas'].append(f'PDF corrompido: {str(e)}')
            resultado['reparavel'] = self._verificar_reparabilidade_pdf(arquivo)
            
        except Exception as e:
            resultado['corrompido'] = True
            resultado['problemas'].append(f'Erro PDF: {str(e)}')
            resultado['reparavel'] = False
        
        return resultado
    
    def _verificar_reparabilidade_zip(self, arquivo: str) -> bool:
        """Verifica se um arquivo baseado em ZIP pode ser reparado."""
        try:
            with zipfile.ZipFile(arquivo, 'r') as zip_file:
                # Tenta listar o conteúdo
                zip_file.namelist()
                return True
        except zipfile.BadZipFile:
            # Verifica se pelo menos parte do arquivo é legível
            try:
                with open(arquivo, 'rb') as f:
                    header = f.read(4)
                    return header == b'PK\x03\x04'  # Assinatura ZIP
            except:
                return False
        except:
            return False
    
    def _verificar_reparabilidade_pdf(self, arquivo: str) -> bool:
        """Verifica se um arquivo PDF pode ser reparado."""
        try:
            with open(arquivo, 'rb') as f:
                header = f.read(8)
                return header.startswith(b'%PDF-')
        except:
            return False
    
    def _calcular_checksum(self, arquivo: str) -> str:
        """Calcula o checksum MD5 do arquivo."""
        hash_md5 = hashlib.md5()
        try:
            with open(arquivo, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_md5.update(chunk)
            return hash_md5.hexdigest()
        except:
            return None

class ReparadorArquivos:
    """Classe para reparar arquivos corrompidos."""
    
    def __init__(self, logger: logging.Logger):
        self.logger = logger
        self.resultados_reparo = {
            'arquivos_reparados': [],
            'arquivos_nao_reparaveis': [],
            'backups_criados': [],
            'erros_reparo': []
        }
    
    def reparar_arquivo(self, arquivo: str, tipo_problema: str, diretorio_backup: str) -> Dict:
        """Tenta reparar um arquivo corrompido."""
        resultado = {
            'arquivo': arquivo,
            'reparado': False,
            'metodo_usado': None,
            'backup_criado': False,
            'arquivo_reparado': None,
            'problemas': []
        }
        
        try:
            # Cria backup antes de tentar reparar
            backup_path = self._criar_backup(arquivo, diretorio_backup)
            if backup_path:
                resultado['backup_criado'] = True
                self.resultados_reparo['backups_criados'].append(backup_path)
            
            extensao = Path(arquivo).suffix.lower()
            
            if extensao in ['.docx', '.xlsx', '.pptx']:
                resultado = self._reparar_office(arquivo, resultado)
            elif extensao == '.pdf':
                resultado = self._reparar_pdf(arquivo, resultado)
            else:
                resultado['problemas'].append('Tipo de arquivo não suportado para reparo')
            
        except Exception as e:
            resultado['problemas'].append(f'Erro durante reparo: {str(e)}')
            self.logger.error(f"Erro ao reparar {arquivo}: {str(e)}")
        
        return resultado
    
    def _criar_backup(self, arquivo: str, diretorio_backup: str) -> Optional[str]:
        """Cria backup do arquivo original."""
        try:
            os.makedirs(diretorio_backup, exist_ok=True)
            nome_arquivo = os.path.basename(arquivo)
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            nome_backup = f"{timestamp}_{nome_arquivo}"
            backup_path = os.path.join(diretorio_backup, nome_backup)
            
            shutil.copy2(arquivo, backup_path)
            self.logger.info(f"Backup criado: {backup_path}")
            return backup_path
            
        except Exception as e:
            self.logger.error(f"Erro ao criar backup de {arquivo}: {str(e)}")
            return None
    
    def _reparar_office(self, arquivo: str, resultado: Dict) -> Dict:
        """Tenta reparar arquivos do Office (DOCX, XLSX, PPTX)."""
        extensao = Path(arquivo).suffix.lower()
        
        # Método 1: Tentar extrair e recriar o ZIP
        try:
            resultado_metodo1 = self._reparar_zip_office(arquivo, extensao)
            if resultado_metodo1['sucesso']:
                resultado.update(resultado_metodo1)
                resultado['metodo_usado'] = 'Reconstrução ZIP'
                return resultado
        except Exception as e:
            resultado['problemas'].append(f'Método ZIP falhou: {str(e)}')
        
        # Método 2: Tentar conversão via texto (apenas para DOCX)
        if extensao == '.docx':
            try:
                resultado_metodo2 = self._reparar_docx_via_texto(arquivo)
                if resultado_metodo2['sucesso']:
                    resultado.update(resultado_metodo2)
                    resultado['metodo_usado'] = 'Conversão via texto'
                    return resultado
            except Exception as e:
                resultado['problemas'].append(f'Método texto falhou: {str(e)}')
        
        # Método 3: Tentar recuperação parcial
        try:
            resultado_metodo3 = self._recuperacao_parcial_office(arquivo, extensao)
            if resultado_metodo3['sucesso']:
                resultado.update(resultado_metodo3)
                resultado['metodo_usado'] = 'Recuperação parcial'
                return resultado
        except Exception as e:
            resultado['problemas'].append(f'Método recuperação parcial falhou: {str(e)}')
        
        return resultado
    
    def _reparar_zip_office(self, arquivo: str, extensao: str) -> Dict:
        """Tenta reparar arquivo Office reconstruindo a estrutura ZIP."""
        resultado = {'sucesso': False, 'arquivo_reparado': None}
        
        with tempfile.TemporaryDirectory() as temp_dir:
            try:
                # Tenta extrair o que for possível
                with zipfile.ZipFile(arquivo, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)
                
                # Cria novo arquivo
                arquivo_reparado = arquivo.replace(extensao, f'_reparado{extensao}')
                
                with zipfile.ZipFile(arquivo_reparado, 'w', zipfile.ZIP_DEFLATED) as new_zip:
                    for root, dirs, files in os.walk(temp_dir):
                        for file in files:
                            file_path = os.path.join(root, file)
                            arc_name = os.path.relpath(file_path, temp_dir)
                            new_zip.write(file_path, arc_name)
                
                # Verifica se o arquivo reparado funciona
                if self._verificar_arquivo_reparado(arquivo_reparado, extensao):
                    resultado['sucesso'] = True
                    resultado['arquivo_reparado'] = arquivo_reparado
                    self.logger.info(f"Arquivo reparado com sucesso: {arquivo_reparado}")
                else:
                    os.remove(arquivo_reparado)
                
            except Exception as e:
                resultado['erro'] = str(e)
        
        return resultado
    
    def _reparar_docx_via_texto(self, arquivo: str) -> Dict:
        """Tenta reparar DOCX extraindo texto e criando novo documento."""
        resultado = {'sucesso': False, 'arquivo_reparado': None}
        
        try:
            # Extrai texto usando docx2txt
            texto = docx2txt.process(arquivo)
            
            if texto and texto.strip():
                # Cria novo documento
                novo_doc = Document()
                
                # Adiciona o texto extraído
                for linha in texto.split('\n'):
                    if linha.strip():
                        novo_doc.add_paragraph(linha)
                
                # Salva o novo documento
                arquivo_reparado = arquivo.replace('.docx', '_reparado_texto.docx')
                novo_doc.save(arquivo_reparado)
                
                resultado['sucesso'] = True
                resultado['arquivo_reparado'] = arquivo_reparado
                self.logger.info(f"DOCX reparado via texto: {arquivo_reparado}")
            
        except Exception as e:
            resultado['erro'] = str(e)
        
        return resultado
    
    def _recuperacao_parcial_office(self, arquivo: str, extensao: str) -> Dict:
        """Tenta recuperação parcial de arquivos Office."""
        resultado = {'sucesso': False, 'arquivo_reparado': None}
        
        try:
            # Cria arquivo de recuperação com estrutura mínima
            arquivo_recuperado = arquivo.replace(extensao, f'_recuperado{extensao}')
            
            if extensao == '.docx':
                doc = Document()
                doc.add_paragraph("ARQUIVO RECUPERADO - Conteúdo original pode ter sido perdido")
                doc.add_paragraph(f"Arquivo original: {os.path.basename(arquivo)}")
                doc.save(arquivo_recuperado)
                
            elif extensao == '.xlsx':
                wb = openpyxl.Workbook()
                ws = wb.active
                ws['A1'] = "ARQUIVO RECUPERADO"
                ws['A2'] = f"Arquivo original: {os.path.basename(arquivo)}"
                wb.save(arquivo_recuperado)
                
            elif extensao == '.pptx':
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = "ARQUIVO RECUPERADO"
                slide.shapes.placeholders[1].text = f"Arquivo original: {os.path.basename(arquivo)}"
                prs.save(arquivo_recuperado)
            
            resultado['sucesso'] = True
            resultado['arquivo_reparado'] = arquivo_recuperado
            self.logger.info(f"Arquivo recuperado parcialmente: {arquivo_recuperado}")
            
        except Exception as e:
            resultado['erro'] = str(e)
        
        return resultado
    
    def _reparar_pdf(self, arquivo: str, resultado: Dict) -> Dict:
        """Tenta reparar arquivos PDF."""
        # Para PDFs, a recuperação é mais limitada
        # Aqui implementaríamos métodos específicos para PDF se necessário
        resultado['problemas'].append('Reparo de PDF não implementado nesta versão')
        return resultado
    
    def _verificar_arquivo_reparado(self, arquivo: str, extensao: str) -> bool:
        """Verifica se o arquivo reparado está funcional."""
        try:
            if extensao == '.docx':
                doc = Document(arquivo)
                return len(doc.paragraphs) >= 0
            elif extensao == '.xlsx':
                wb = openpyxl.load_workbook(arquivo)
                wb.close()
                return True
            elif extensao == '.pptx':
                prs = Presentation(arquivo)
                return len(prs.slides) >= 0
            return False
        except:
            return False

class ReparadorAvancado:
    """Classe principal para detecção e reparo avançado de arquivos."""
    
    def __init__(self, diretorio_base: str):
        self.diretorio_base = diretorio_base
        self.diretorio_logs = os.path.join(diretorio_base, "logs")
        self.diretorio_backup = os.path.join(diretorio_base, "backups_reparo")
        self.diretorio_reparados = os.path.join(diretorio_base, "arquivos_reparados")
        
        # Configura logger
        self.logger = LoggerConfig.setup_logger(self.diretorio_logs)
        
        # Inicializa componentes
        self.detector = DetectorCorrupcao(self.logger)
        self.reparador = ReparadorArquivos(self.logger)
        
        # Cria diretórios necessários
        for diretorio in [self.diretorio_backup, self.diretorio_reparados]:
            os.makedirs(diretorio, exist_ok=True)
    
    def executar_analise_completa(self) -> Dict:
        """Executa análise completa de detecção e reparo."""
        self.logger.info("Iniciando análise avançada de detecção e reparo...")
        
        # Encontra todos os arquivos
        arquivos = self._encontrar_arquivos()
        
        resultados = {
            'timestamp': datetime.now().strftime("%Y%m%d_%H%M%S"),
            'total_arquivos': len(arquivos),
            'arquivos_analisados': 0,
            'arquivos_corrompidos': 0,
            'arquivos_reparados': 0,
            'arquivos_nao_reparaveis': 0,
            'detalhes_arquivos': [],
            'estatisticas': {},
            'erros': []
        }
        
        # Analisa cada arquivo
        for arquivo in arquivos:
            try:
                resultado_arquivo = self._analisar_arquivo(arquivo)
                resultados['detalhes_arquivos'].append(resultado_arquivo)
                resultados['arquivos_analisados'] += 1
                
                if resultado_arquivo.get('corrompido', False):
                    resultados['arquivos_corrompidos'] += 1
                    
                    # Tenta reparar se possível
                    if resultado_arquivo.get('reparavel', False):
                        resultado_reparo = self.reparador.reparar_arquivo(
                            arquivo, 
                            resultado_arquivo.get('problemas', []),
                            self.diretorio_backup
                        )
                        
                        resultado_arquivo['reparo'] = resultado_reparo
                        
                        if resultado_reparo.get('reparado', False):
                            resultados['arquivos_reparados'] += 1
                            # Move arquivo reparado para diretório específico
                            self._mover_arquivo_reparado(resultado_reparo['arquivo_reparado'])
                        else:
                            resultados['arquivos_nao_reparaveis'] += 1
                
            except Exception as e:
                erro = f"Erro ao analisar {arquivo}: {str(e)}"
                resultados['erros'].append(erro)
                self.logger.error(erro)
        
        # Gera estatísticas
        resultados['estatisticas'] = self._gerar_estatisticas(resultados)
        
        # Salva relatório
        self._salvar_relatorio(resultados)
        
        # Exibe resumo
        self._exibir_resumo(resultados)
        
        return resultados
    
    def _encontrar_arquivos(self) -> List[str]:
        """Encontra todos os arquivos para análise."""
        arquivos = []
        extensoes_suportadas = {'.docx', '.xlsx', '.pptx', '.pdf'}
        
        for root, dirs, files in os.walk(self.diretorio_base):
            # Pula diretórios de sistema
            if any(skip in root for skip in ['logs', 'backups', 'arquivos_reparados', '__pycache__']):
                continue
                
            for file in files:
                if Path(file).suffix.lower() in extensoes_suportadas:
                    arquivos.append(os.path.join(root, file))
        
        return arquivos
    
    def _analisar_arquivo(self, arquivo: str) -> Dict:
        """Analisa um arquivo específico."""
        resultado = {
            'arquivo': arquivo,
            'nome': os.path.basename(arquivo),
            'extensao': Path(arquivo).suffix.lower(),
            'tamanho': os.path.getsize(arquivo),
            'corrompido': False,
            'reparavel': False,
            'problemas': [],
            'detalhes': {}
        }
        
        # Verificação básica
        verificacao_basica = self.detector.verificar_integridade_basica(arquivo)
        resultado['detalhes']['verificacao_basica'] = verificacao_basica
        
        if verificacao_basica['problemas']:
            resultado['corrompido'] = True
            resultado['problemas'].extend(verificacao_basica['problemas'])
        
        # Verificação específica por tipo
        extensao = resultado['extensao']
        
        if extensao == '.docx':
            verificacao_especifica = self.detector.verificar_docx(arquivo)
        elif extensao == '.xlsx':
            verificacao_especifica = self.detector.verificar_xlsx(arquivo)
        elif extensao == '.pptx':
            verificacao_especifica = self.detector.verificar_pptx(arquivo)
        elif extensao == '.pdf':
            verificacao_especifica = self.detector.verificar_pdf(arquivo)
        else:
            verificacao_especifica = {'tipo': 'DESCONHECIDO', 'corrompido': False, 'problemas': []}
        
        resultado['detalhes']['verificacao_especifica'] = verificacao_especifica
        
        if verificacao_especifica.get('corrompido', False):
            resultado['corrompido'] = True
            resultado['problemas'].extend(verificacao_especifica.get('problemas', []))
            resultado['reparavel'] = verificacao_especifica.get('reparavel', False)
        
        return resultado
    
    def _mover_arquivo_reparado(self, arquivo_reparado: str):
        """Move arquivo reparado para diretório específico."""
        try:
            nome_arquivo = os.path.basename(arquivo_reparado)
            destino = os.path.join(self.diretorio_reparados, nome_arquivo)
            shutil.move(arquivo_reparado, destino)
            self.logger.info(f"Arquivo reparado movido para: {destino}")
        except Exception as e:
            self.logger.error(f"Erro ao mover arquivo reparado: {str(e)}")
    
    def _gerar_estatisticas(self, resultados: Dict) -> Dict:
        """Gera estatísticas detalhadas."""
        estatisticas = {
            'taxa_corrupcao': 0,
            'taxa_reparo': 0,
            'tipos_arquivo': {},
            'problemas_comuns': {},
            'metodos_reparo_usados': {}
        }
        
        if resultados['total_arquivos'] > 0:
            estatisticas['taxa_corrupcao'] = (resultados['arquivos_corrompidos'] / resultados['total_arquivos']) * 100
        
        if resultados['arquivos_corrompidos'] > 0:
            estatisticas['taxa_reparo'] = (resultados['arquivos_reparados'] / resultados['arquivos_corrompidos']) * 100
        
        # Analisa tipos de arquivo
        for detalhe in resultados['detalhes_arquivos']:
            extensao = detalhe['extensao']
            if extensao not in estatisticas['tipos_arquivo']:
                estatisticas['tipos_arquivo'][extensao] = {'total': 0, 'corrompidos': 0, 'reparados': 0}
            
            estatisticas['tipos_arquivo'][extensao]['total'] += 1
            
            if detalhe['corrompido']:
                estatisticas['tipos_arquivo'][extensao]['corrompidos'] += 1
                
                if detalhe.get('reparo', {}).get('reparado', False):
                    estatisticas['tipos_arquivo'][extensao]['reparados'] += 1
                    
                    # Conta métodos de reparo
                    metodo = detalhe['reparo'].get('metodo_usado', 'Desconhecido')
                    estatisticas['metodos_reparo_usados'][metodo] = estatisticas['metodos_reparo_usados'].get(metodo, 0) + 1
            
            # Conta problemas comuns
            for problema in detalhe['problemas']:
                estatisticas['problemas_comuns'][problema] = estatisticas['problemas_comuns'].get(problema, 0) + 1
        
        return estatisticas
    
    def _salvar_relatorio(self, resultados: Dict):
        """Salva relatório detalhado em JSON."""
        timestamp = resultados['timestamp']
        arquivo_relatorio = os.path.join(self.diretorio_logs, f"relatorio_reparo_avancado_{timestamp}.json")
        
        try:
            with open(arquivo_relatorio, 'w', encoding='utf-8') as f:
                json.dump(resultados, f, indent=2, ensure_ascii=False)
            
            self.logger.info(f"Relatório salvo em: {arquivo_relatorio}")
            
        except Exception as e:
            self.logger.error(f"Erro ao salvar relatório: {str(e)}")
    
    def _exibir_resumo(self, resultados: Dict):
        """Exibe resumo dos resultados."""
        print("\n" + "="*60)
        print("RELATÓRIO FINAL - DETECÇÃO E REPARO AVANÇADO")
        print("="*60)
        print(f"Total de arquivos analisados: {resultados['arquivos_analisados']}")
        print(f"Arquivos corrompidos encontrados: {resultados['arquivos_corrompidos']}")
        print(f"Arquivos reparados com sucesso: {resultados['arquivos_reparados']}")
        print(f"Arquivos não reparáveis: {resultados['arquivos_nao_reparaveis']}")
        print(f"Taxa de corrupção: {resultados['estatisticas']['taxa_corrupcao']:.2f}%")
        print(f"Taxa de reparo: {resultados['estatisticas']['taxa_reparo']:.2f}%")
        
        if resultados['estatisticas']['tipos_arquivo']:
            print("\nEstatísticas por tipo de arquivo:")
            for tipo, stats in resultados['estatisticas']['tipos_arquivo'].items():
                print(f"  {tipo}: {stats['total']} total, {stats['corrompidos']} corrompidos, {stats['reparados']} reparados")
        
        if resultados['estatisticas']['metodos_reparo_usados']:
            print("\nMétodos de reparo utilizados:")
            for metodo, count in resultados['estatisticas']['metodos_reparo_usados'].items():
                print(f"  {metodo}: {count} vezes")
        
        print("="*60)

def main():
    """Função principal."""
    # Diretório base (pode ser alterado conforme necessário)
    diretorio_base = "/Users/marcosdaniels/Downloads/5. FULL SALES 2/Consolidado"
    
    if not os.path.exists(diretorio_base):
        print(f"Erro: Diretório {diretorio_base} não encontrado!")
        return
    
    # Cria e executa o reparador avançado
    reparador = ReparadorAvancado(diretorio_base)
    resultados = reparador.executar_analise_completa()
    
    print(f"\nProcesso concluído! Verifique os logs em: {reparador.diretorio_logs}")
    if resultados['arquivos_reparados'] > 0:
        print(f"Arquivos reparados salvos em: {reparador.diretorio_reparados}")

if __name__ == "__main__":
    main()